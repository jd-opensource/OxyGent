import re
import pdfplumber
import requests
import tempfile
import os
import mimetypes
import chardet
from pptx import Presentation
from openai import OpenAI
from oxygent import oxy
from pydantic import Field
from typing import Optional, Dict, Union, List
from urllib.parse import urlparse
from pydantic.fields import FieldInfo


doc_qa_tools = oxy.FunctionHub(name="document_qa_tools", timeout=600)


def download_file(url: str) -> Optional[str]:
    """Download file to temporary location and return local path"""
    try:
        # Extract file extension from URL
        path = urlparse(url).path
        file_ext = os.path.splitext(path)[1] if path else None

        response = requests.get(url, stream=True)
        response.raise_for_status()

        # Create temporary file with original extension if available
        suffix = file_ext if file_ext else '.tmp'
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)

        # Write file in chunks
        for chunk in response.iter_content(chunk_size=8192):
            if chunk:
                temp_file.write(chunk)

        temp_file.close()
        return temp_file.name

    except Exception as e:
        print(f"File download failed: {e}")
        return None


def determine_file_type(file_path: str, content: bytes = None) -> str:
    """Determine file type, prioritizing extension with fallback to content detection"""
    # Try based on extension
    ext = os.path.splitext(file_path)[-1][1:].lower()
    if ext in ['pdf', 'txt', 'ppt', 'pptx']:
        return ext

    # Try content-based detection (for downloaded URLs)
    if content:
        mime = mimetypes.guess_type(file_path)[0]
        if mime:
            if 'pdf' in mime:
                return 'pdf'
            elif 'text' in mime:
                return 'txt'
            elif 'powerpoint' in mime or 'presentation' in mime:
                return 'ppt'

    # Default fallback
    return 'txt'


def normalize_context_range(context_range: Union[int, str]) -> int:
    """Ensure context range is an integer"""
    try:
        if isinstance(context_range, str):
            return int(context_range.strip())
        return int(context_range)
    except (ValueError, TypeError):
        return 300


def retrieve_content(file_path: str, search_text: str, context_range: Union[int, str] = 300,
                     file_type: Optional[str] = None) -> Dict[str, str]:
    """Retrieve document content, return location info and extracted text snippet"""
    result = {"position": "", "content": ""}
    local_file = None
    temp_file = None

    try:
        # If URL, download file
        if file_path.startswith(('http://', 'https://')):
            print(f"Downloading URL content: {file_path}")
            temp_file = download_file(file_path)
            if not temp_file:
                return {"position": "Error", "content": "File download failed"}
            local_file = temp_file
            print(f"Downloaded to temporary location: {local_file}")
        else:
            local_file = file_path

        # Verify file exists
        if not os.path.exists(local_file):
            return {"position": "Error", "content": f"File not found: {local_file}"}

        # Determine file type
        if not file_type:
            with open(local_file, 'rb') as f:
                content = f.read(1024)  # Read first 1KB for type detection
                file_type = determine_file_type(local_file, content)

        # Normalize context range
        context_range = normalize_context_range(context_range)
        print(f"File type: {file_type}, Context range: {context_range}")

        # Process different file types
        if file_type == 'pdf':
            with pdfplumber.open(local_file) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    text = page.extract_text() or ""
                    if re.search(search_text, text, re.IGNORECASE):
                        result["position"] = f"Page {page_num}"
                        result["content"] = text
                        return result

        elif file_type == 'txt':
            # Auto-detect encoding
            with open(local_file, 'rb') as f:
                raw_data = f.read()
                encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
                full_text = raw_data.decode(encoding, errors='replace')

            match = re.search(search_text, full_text, re.IGNORECASE)
            if match:
                start_pos = max(0, match.start() - context_range)
                end_pos = min(len(full_text), match.end() + context_range)
                context = full_text[start_pos:end_pos]
                result["position"] = f"Character {match.start()}-{match.end()}"
                result["content"] = context
                return result

        elif file_type in ['ppt', 'pptx']:
            prs = Presentation(local_file)
            for slide_num, slide in enumerate(prs.slides, start=1):
                text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                if text and re.search(search_text, text, re.IGNORECASE):
                    result["position"] = f"Slide {slide_num}"
                    result["content"] = text
                    return result

        result["content"] = f"No content found containing keyword '{search_text}'"

    except Exception as e:
        result["content"] = f"Processing error: {str(e)}"
        return result

    finally:
        # Clean up temporary files
        if temp_file and os.path.exists(temp_file):
            os.unlink(temp_file)
            print(f"Cleaned temporary file: {temp_file}")

    return result


def query_llm(content: str, question: str, model: str = "gpt-4o-0806") -> str:
    """Queries LLM with content to answer question"""
    prompt = f"""
    [Reference Content]
    {content}

    [Question]
    {question}

    Note: Answer strictly based on reference content. Do not add external knowledge.
    """
    # Placeholder for actual API call
    client4o = OpenAI(
        base_url=os.getenv('OPEN_AI_URL'),
        api_key=os.getenv('OPEN_AI_KEY'),
    )
    response = client4o.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content


# @doc_qa_tools.tool(description="End-to-end document intelligence system that: "
#                                "1. SEARCHES for target text patterns in PDF/TXT/PPT files/url "
#                                "2. EXTRACTS precise context around matches "
#                                "3. ANSWERS natural language questions using AI reasoning "
#                                "Solves: Rapid information retrieval from business documents without manual review")
def document_qa_api(
        file_path: str = Field(description="Path to the document file",
                               examples=["/data/report.pdf", "/notes/chapter.txt"]),
        search_text: str = Field(description="Regular expression pattern to search in the document",
                                 examples=["market analysis", "Q[1-4] results"]),
        question: str = Field(description="Question to ask the LLM based on the content",
                              examples=["What was the growth rate?"]),
        context_range: Union[int, str] = Field(
            description="Number of characters to extract around the match (TXT only). Default: 300", default=300),
) -> str:
    """
    Retrieves document content by locating search text and extracting surrounding context
    Returns dict with keys: position (location info), content (extracted text)
    """
    if isinstance(context_range, FieldInfo):
        context_range = context_range.default
    retrieved_ = retrieve_content(file_path, search_text, context_range)
    print(f"Position: {retrieved_['position']}\nLength: {len(retrieved_['content'])} character")
    return query_llm(retrieved_["content"], question=question)


if __name__ == "__main__":
    from dotenv import load_dotenv
    from pathlib import Path

    env_path = Path('../examples/gaia/') / '.env'
    load_dotenv(dotenv_path=env_path, verbose=True)
    DOC_PATH = "pdfcoffee.com_offcial-blackbook-price-guide-to-united-states-paper-money-2006-pdf-free.pdf"
    SEARCH_TEXT = r"extra.*pieces"
    QUESTION = "what is offering by Carhsle Development ？"
    print(document_qa_api(DOC_PATH, SEARCH_TEXT, QUESTION))

    #####
    DOC_PATH = "https://ia800303.us.archive.org/34/items/AHandbookForLearningTogether/F-Arato-A-Varga-A-Handbook-for-learning-together-University-of-Pecs-Pecs-2015.pdf"
    SEARCH_TEXT = r"in the past two periods"
    QUESTION = "how much minutes the author have no idea what they talked.？"

    print(document_qa_api(DOC_PATH, SEARCH_TEXT, QUESTION))
